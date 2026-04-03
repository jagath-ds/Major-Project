"""
Document Loader — Sprint 1
Reads files from docs/ folder and converts them into page dicts.
Supported: PDF, DOCX, TXT, HTML, PPTX
"""

from __future__ import annotations

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import logging
from pathlib import Path
from typing import List, Dict, Optional

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".html", ".htm", ".pptx"}


class DocumentLoader:
    def __init__(self, docs_dir: str = "../docs", recursive: bool = True):
        self.docs_dir  = Path(docs_dir)
        self.recursive = recursive
        if not self.docs_dir.exists():
            self.docs_dir.mkdir(parents=True)
            logger.info(f"Created docs folder: {self.docs_dir.resolve()}")
        else:
            logger.info(f"Docs folder: {self.docs_dir.resolve()}")

    def load_all(self) -> List[Dict]:
        files = self._find_files()
        if not files:
            logger.warning(f"No supported documents found in '{self.docs_dir}'.")
            return []
        documents = []
        for filepath in files:
            doc = self._load_file(filepath)
            if doc:
                documents.append(doc)
        logger.info(f"Loaded {len(documents)} documents.")
        return documents

    def load_file(self, filepath: str) -> Optional[Dict]:
        return self._load_file(Path(filepath))

    def _find_files(self) -> List[Path]:
        pattern = "**/*" if self.recursive else "*"
        files = [
            p for p in self.docs_dir.glob(pattern)
            if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS
        ]
        files.sort()
        logger.info(f"Found {len(files)} supported files.")
        return files

    def _load_file(self, filepath: Path) -> Optional[Dict]:
        ext = filepath.suffix.lower()
        logger.info(f"Loading: {filepath.name} [{ext}]")
        try:
            if ext == ".pdf":
                pages = self._load_pdf(filepath)
            elif ext == ".docx":
                pages = self._load_docx(filepath)
            elif ext == ".txt":
                pages = self._load_txt(filepath)
            elif ext in (".html", ".htm"):
                pages = self._load_html(filepath)
            elif ext == ".pptx":
                pages = self._load_pptx(filepath)
            else:
                return None
            if not pages:
                logger.warning(f"No text extracted from: {filepath.name}")
                return None
            return {
                "doc_id":      self._make_doc_id(filepath),
                "source_path": str(filepath),
                "file_type":   ext.lstrip("."),
                "pages":       pages,
            }
        except Exception as e:
            logger.error(f"Failed to load {filepath.name}: {e}")
            return None

    def _load_pdf(self, filepath: Path) -> List[Dict]:
        try:
            import pdfplumber
        except ImportError:
            raise ImportError("Run: pip install pdfplumber")
        pages = []
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = (page.extract_text() or "").strip()
                if text:
                    pages.append({"page": i, "text": text})
        return pages

    def _load_docx(self, filepath: Path) -> List[Dict]:
        try:
            from docx import Document
        except ImportError:
            raise ImportError("Run: pip install python-docx")
        doc      = Document(filepath)
        all_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return self._split_into_pages(all_text)

    def _load_txt(self, filepath: Path) -> List[Dict]:
        text = filepath.read_text(encoding="utf-8", errors="ignore").strip()
        return self._split_into_pages(text)

    def _load_html(self, filepath: Path) -> List[Dict]:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("Run: pip install beautifulsoup4")
        html = filepath.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        return self._split_into_pages(soup.get_text(separator="\n").strip())

    def _load_pptx(self, filepath: Path) -> List[Dict]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Run: pip install python-pptx")
        prs   = Presentation(filepath)
        pages = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text.strip() for shape in slide.shapes
                     if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                pages.append({"page": i, "text": "\n".join(texts)})
        return pages

    @staticmethod
    def _split_into_pages(text: str, chunk_size: int = 3000) -> List[Dict]:
        pages = []
        for i, start in enumerate(range(0, len(text), chunk_size), 1):
            chunk = text[start: start + chunk_size].strip()
            if chunk:
                pages.append({"page": i, "text": chunk})
        return pages

    @staticmethod
    def _make_doc_id(filepath: Path) -> str:
        stem = filepath.stem.lower().replace(" ", "_").replace("-", "_")
        return "".join(c for c in stem if c.isalnum() or c == "_")